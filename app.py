from flask import Flask, render_template, request, jsonify, send_from_directory
import os
import pyttsx3
from pptx import Presentation
from zipfile import ZipFile

# Initialize the Flask app
app = Flask(__name__)

# Folder to store uploaded PPT files and generated audio
UPLOAD_FOLDER = 'uploads'
AUDIO_FOLDER = 'audio_files'

# Ensure the folders exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(AUDIO_FOLDER, exist_ok=True)

# Configure Flask to serve static files
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['AUDIO_FOLDER'] = AUDIO_FOLDER

# Function to convert text to speech using pyttsx3
def text_to_speech(text, filename):
    engine = pyttsx3.init()
    audio_path = os.path.join(AUDIO_FOLDER, filename)
    engine.save_to_file(text, audio_path)
    engine.runAndWait()
    return audio_path

# Function to extract text from a PowerPoint file
def extract_text_from_ppt(ppt_file):
    try:
        prs = Presentation(ppt_file)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        return str(e)

# Route to display the HTML form
@app.route('/')
def index():
    return render_template('index.html')

# Route to handle PPT file upload and process it
@app.route('/upload', methods=['POST'])
def upload_ppt():
    if 'ppt-file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    ppt_file = request.files['ppt-file']
    if ppt_file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    # Check if the file is either .pptx or .ppt
    if ppt_file and (ppt_file.filename.endswith('.pptx') or ppt_file.filename.endswith('.ppt')):
        # Save the uploaded file
        ppt_file_path = os.path.join(app.config['UPLOAD_FOLDER'], ppt_file.filename)
        ppt_file.save(ppt_file_path)

        # Extract text from the file
        text = extract_text_from_ppt(ppt_file_path)
        if not text:
            return jsonify({'error': 'Failed to extract text from PowerPoint file.'}), 400
        
        # Convert text to speech and save it as an audio file
        audio_filename = f"{ppt_file.filename.split('.')[0]}.mp3"
        audio_path = text_to_speech(text, audio_filename)

        # Return the audio file URL to the client
        return jsonify({"audioUrl": f"/audio/{audio_filename}"}), 200

    return jsonify({'error': 'Invalid file format. Please upload a .ppt or .pptx file.'}), 400

# Route to serve the generated audio file
@app.route('/audio/<filename>')
def serve_audio(filename):
    return send_from_directory(AUDIO_FOLDER, filename)

# Run the app
if __name__ == "__main__":
    app.run(debug=True)
